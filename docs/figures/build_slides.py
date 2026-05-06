"""Generate the technical-audience PowerPoint deck.

Run AFTER ``build_figures.py`` (which produces the PNGs this deck embeds):

    PYTHONPATH=src .venv/bin/python docs/figures/build_figures.py
    PYTHONPATH=src .venv/bin/python docs/figures/build_slides.py

Output: ``docs/architecture_slides.pptx`` (16:9, ~14 slides).

Audience
--------
Technical — fresh engineer or operator joining the project. No
sales pitch, no decorative visuals; every slide either embeds one
of the architecture diagrams or surfaces concrete facts that
support it.

Design
------
- 16:9 widescreen (Microsoft default for technical decks).
- Single-color title bar, generous whitespace, no logos.
- Helvetica/Arial fallback chain (matches the figures).
- One topic per slide; bullets ≤ 5; no walls of text.
- Bullet body text 18pt — readable from the back of a 6m room.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

FIG_DIR = Path(__file__).parent
OUT_PATH = FIG_DIR.parent / "architecture_slides.pptx"

# Color palette — kept in sync with the figure palette
TITLE_BAR = RGBColor(0x2F, 0x4F, 0x6F)
TITLE_TXT = RGBColor(0xFF, 0xFF, 0xFF)
BODY_TXT = RGBColor(0x22, 0x22, 0x22)
ACCENT = RGBColor(0xD1, 0x68, 0x37)
SUBTLE = RGBColor(0x66, 0x66, 0x66)

FONT = "Helvetica"
FALLBACK_FONT = "Arial"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


def _new_slide(prs: Presentation):
    """Blank layout, custom title bar drawn manually."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _add_title_bar(slide, title: str, subtitle: str | None = None) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        SLIDE_W,
        Inches(0.85),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TITLE_BAR
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.18)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = title
    run = p.runs[0]
    run.font.name = FONT
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = TITLE_TXT
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        run2 = p2.runs[0]
        run2.font.name = FONT
        run2.font.size = Pt(13)
        run2.font.color.rgb = RGBColor(0xCC, 0xD8, 0xE5)


def _add_text_block(slide, left, top, width, height, text, *, fontsize=18, bold=False, color=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.name = FONT
    run.font.size = Pt(fontsize)
    run.font.bold = bold
    run.font.color.rgb = color or BODY_TXT


def _add_bullets(slide, left, top, width, height, items, *, fontsize=18, line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + text
        p.line_spacing = line_spacing
        run = p.runs[0]
        run.font.name = FONT
        run.font.size = Pt(fontsize)
        run.font.color.rgb = BODY_TXT


def _add_image(slide, png_path, left, top, width=None, height=None):
    if width and height:
        slide.shapes.add_picture(str(png_path), left, top, width=width, height=height)
    elif width:
        slide.shapes.add_picture(str(png_path), left, top, width=width)
    elif height:
        slide.shapes.add_picture(str(png_path), left, top, height=height)
    else:
        slide.shapes.add_picture(str(png_path), left, top)


def _add_footer(slide, page_num, total):
    ftr = slide.shapes.add_textbox(Inches(11.5), Inches(7.05), Inches(1.7), Inches(0.4))
    tf = ftr.text_frame
    p = tf.paragraphs[0]
    p.text = f"{page_num} / {total}"
    run = p.runs[0]
    run.font.name = FONT
    run.font.size = Pt(10)
    run.font.color.rgb = SUBTLE


# ---------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------


def slide_title(prs):
    s = _new_slide(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.2), SLIDE_W, Inches(3.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TITLE_BAR
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.7)
    tf.margin_top = Inches(0.4)
    p = tf.paragraphs[0]
    p.text = "APECx MCP Integration"
    r = p.runs[0]
    r.font.name = FONT
    r.font.size = Pt(48)
    r.font.bold = True
    r.font.color.rgb = TITLE_TXT
    p2 = tf.add_paragraph()
    p2.text = "End-to-End Architecture"
    r2 = p2.runs[0]
    r2.font.name = FONT
    r2.font.size = Pt(32)
    r2.font.color.rgb = RGBColor(0xCC, 0xD8, 0xE5)
    p3 = tf.add_paragraph()
    p3.text = "Synthesis pipeline · 23 MCP tools · 6 ontologies · 504 unit tests"
    r3 = p3.runs[0]
    r3.font.name = FONT
    r3.font.size = Pt(18)
    r3.font.color.rgb = RGBColor(0xCC, 0xD8, 0xE5)
    _add_text_block(
        s,
        Inches(0.7),
        Inches(6.4),
        Inches(11.0),
        Inches(0.8),
        "Generated 2026-05-05  ·  apecx-mcp-integration repo  ·  branch day2-rag-synthesis-agent",
        fontsize=14,
        color=SUBTLE,
    )


def slide_agenda(prs, page, total):
    s = _new_slide(prs)
    _add_title_bar(s, "Agenda", "What this deck covers")
    items = [
        "Two lifecycles — backend (offline, periodic) vs. user-facing (online, per-query)",
        "Backend harmonization pipeline — harvester + dictionary builder",
        "User-facing query workflow — MCP entry → synthesis → markdown",
        "Three-tier runtime topology (Tiers 1, 2, 4)",
        "Synthesis pipeline detail — 5 retrieval branches + LLM",
        "Trigger-cascade primitive (added 2026-05-05 to nanobrain)",
        "Mapping & resolution strategy — fast / ancestor / slow / miss",
        "Ontologies, MCP tool surface, test surface",
        "Data quality — CI-enforced accuracy floors & harmonization statistics",
        "Things that will surprise you (silent-failure shapes)",
    ]
    _add_bullets(s, Inches(0.8), Inches(1.3), Inches(12.0), Inches(5.5), items, fontsize=20)
    _add_footer(s, page, total)


def slide_purpose(prs, page, total):
    s = _new_slide(prs)
    _add_title_bar(s, "System Purpose", "One-paragraph definition")
    _add_text_block(
        s,
        Inches(0.8),
        Inches(1.4),
        Inches(11.7),
        Inches(2.5),
        (
            "APECx MCP integration takes free-text scientist questions about viral "
            "pathogens, vaccines, genes, and genomes and turns them into grounded "
            "Markdown answers with inline citations."
        ),
        fontsize=22,
        bold=True,
    )
    _add_text_block(
        s,
        Inches(0.8),
        Inches(3.5),
        Inches(11.7),
        Inches(3.0),
        (
            "It does this by fanning out across local FAISS-indexed knowledge (domain RAG), "
            "local VIOLIN/BV-BRC tabular data (substring lookup), live PubMed publications, "
            "and the APECx Globus Search index of harvested records — then driving one LLM "
            "call to weave the retrieved evidence into a structured response.\n\n"
            "The system is exposed via the Model Context Protocol (MCP), so it appears as a "
            "tool surface inside Claude Desktop and any MCP-compatible client."
        ),
        fontsize=17,
    )
    _add_footer(s, page, total)


def slide_two_lifecycles(prs, page, total):
    """The orienting slide: this system has two distinct lifecycles
    and confusing them is the #1 onboarding mistake."""
    s = _new_slide(prs)
    _add_title_bar(
        s,
        "Two Lifecycles, One Bridge",
        "Backend vs. user-facing — the central organizing distinction",
    )

    # Left column: Backend
    _add_text_block(
        s,
        Inches(0.7),
        Inches(1.3),
        Inches(6.0),
        Inches(0.6),
        "BACKEND HARMONIZATION",
        fontsize=18,
        bold=True,
        color=ACCENT,
    )
    backend_items = [
        "Run periodically (e.g. monthly)",
        "Two pipelines: harvester (apecx-harvesters) + dictionary builder (this repo)",
        "Writes Globus Search index + apecx_synonym_dict.sqlite",
        "Synchronous, batch, NOT on the per-query hot path",
        "Sources: PubMed, PDB, DataCite, Crossref, OpenAlex, bioRxiv, EMDB, DOI, VIOLIN, BV-BRC, NCBI taxdump",
    ]
    _add_bullets(
        s,
        Inches(0.7),
        Inches(1.95),
        Inches(6.0),
        Inches(4.5),
        backend_items,
        fontsize=14,
        line_spacing=1.25,
    )

    # Right column: User-facing
    _add_text_block(
        s,
        Inches(7.0),
        Inches(1.3),
        Inches(6.0),
        Inches(0.6),
        "USER-FACING WORKFLOW",
        fontsize=18,
        bold=True,
        color=RGBColor(0x30, 0x60, 0xA0),
    )
    user_items = [
        "Per scientist query (~70s wall-clock on Ollama)",
        "MCP server → synthesize_query → 4 retrieval branches → LLM",
        "Reads the artifacts the backend wrote — never writes back",
        "Real-time, async, fail-soft branch failures",
        "Globus + FAISS + VIOLIN/BV-BRC + PubMed + (synonym dict for fast-path resolution)",
    ]
    _add_bullets(
        s,
        Inches(7.0),
        Inches(1.95),
        Inches(6.0),
        Inches(4.5),
        user_items,
        fontsize=14,
        line_spacing=1.25,
    )

    # Bottom: bridge sentence
    _add_text_block(
        s,
        Inches(0.7),
        Inches(6.3),
        Inches(12.0),
        Inches(0.7),
        "Bridge: every artifact the backend writes, the user-facing workflow reads — and never writes to.",
        fontsize=15,
        bold=True,
        color=BODY_TXT,
    )
    _add_footer(s, page, total)


def slide_backend_pipeline(prs, page, total):
    s = _new_slide(prs)
    _add_title_bar(
        s,
        "Backend Harmonization Pipeline",
        "Offline · run periodically · harvester (apecx-harvesters) + dictionary builder (this repo)",
    )
    _add_image(
        s, FIG_DIR / "10_backend_harmonization.png", Inches(0.4), Inches(1.0), width=Inches(12.6)
    )
    _add_footer(s, page, total)


def slide_user_facing_pipeline(prs, page, total):
    s = _new_slide(prs)
    _add_title_bar(
        s,
        "User-Facing Workflow",
        "Online · per-query · ~70s wall-clock · 3 bands (MCP entry → synthesis → artifacts read)",
    )
    _add_image(
        s, FIG_DIR / "11_user_facing_workflow.png", Inches(0.4), Inches(1.0), width=Inches(12.6)
    )
    _add_footer(s, page, total)


def slide_topology(prs, page, total):
    s = _new_slide(prs)
    _add_title_bar(
        s,
        "Three-Tier Runtime Topology",
        "Tier 1 (this repo) · Tier 2 (Control Plane) · Tier 4 (Executors)",
    )
    _add_image(
        s, FIG_DIR / "01_three_tier_topology.png", Inches(0.4), Inches(1.0), width=Inches(12.6)
    )
    _add_footer(s, page, total)


def slide_synthesis(prs, page, total):
    s = _new_slide(prs)
    _add_title_bar(
        s,
        "Synthesis Pipeline",
        "5 concurrent retrieval branches → 1 LLM round-trip → grounded Markdown",
    )
    _add_image(
        s, FIG_DIR / "02_synthesis_pipeline.png", Inches(0.4), Inches(1.0), width=Inches(12.6)
    )
    _add_footer(s, page, total)


def slide_invocation(prs, page, total):
    s = _new_slide(prs)
    _add_title_bar(s, "Three Invocation Paths", "Pick one based on the calling context")
    _add_image(s, FIG_DIR / "03_invocation_paths.png", Inches(0.4), Inches(1.0), width=Inches(12.6))
    _add_footer(s, page, total)


def slide_trigger_cascade(prs, page, total):
    s = _new_slide(prs)
    _add_title_bar(
        s,
        "Workflow.wait_for_cascade",
        "Added to nanobrain 2026-05-05 — drains the trigger executor's task set with a settle window",
    )
    _add_image(
        s, FIG_DIR / "08_trigger_cascade_timeline.png", Inches(0.4), Inches(1.0), width=Inches(12.6)
    )
    _add_footer(s, page, total)


def slide_resolution(prs, page, total):
    s = _new_slide(prs)
    _add_title_bar(
        s,
        "Mapping & Resolution Strategy",
        "fast → ancestor → slow → miss; every result carries (path, confidence)",
    )
    _add_image(
        s, FIG_DIR / "04_resolution_decision_tree.png", Inches(2.0), Inches(1.0), width=Inches(9.5)
    )
    _add_footer(s, page, total)


def slide_ontologies(prs, page, total):
    s = _new_slide(prs)
    _add_title_bar(s, "Ontologies", "Authoritative sources tracked in DictionaryEntry.ontology")
    _add_image(
        s, FIG_DIR / "06_ontologies_coverage.png", Inches(0.4), Inches(1.0), width=Inches(12.6)
    )
    _add_footer(s, page, total)


def slide_mcp_tools(prs, page, total):
    s = _new_slide(prs)
    _add_title_bar(s, "MCP Tool Surface", "23 tools across 8 groups")
    _add_image(
        s, FIG_DIR / "07_mcp_tool_distribution.png", Inches(0.7), Inches(1.0), width=Inches(11.9)
    )
    _add_footer(s, page, total)


def slide_tests(prs, page, total):
    s = _new_slide(prs)
    _add_title_bar(
        s,
        "Test Surface",
        "504 unit tests · 7 workflow YAML · 1 cascade runtime — auto-skip when external deps missing",
    )
    _add_image(s, FIG_DIR / "05_test_surface.png", Inches(0.4), Inches(1.0), width=Inches(12.6))
    _add_footer(s, page, total)


def slide_accuracy_methodology(prs, page, total):
    s = _new_slide(prs)
    _add_title_bar(
        s,
        "Data Quality Assessment",
        "Three-test pattern · CI-enforced floors · live-OLS gating, no mocks",
    )
    items = [
        "AccuracyMetrics — recall + precision + F1 from confusion-matrix counts (correct / incorrect / ground-truth total)",
        "Three-test pattern: slice baseline (60-row deterministic sample, fast feedback) → "
        "full corpus (all 13,238 rows, gated by APECX_SYNONYM_DICT_FULL_CORPUS=1) → "
        "probe-batch sampling (50 / 300 spot-checks at the boundary)",
        "Live OLS gating — APECX_SYNONYM_DICT_LIVE_OLS=1 required; tests auto-skip when "
        "the resolver is unreachable rather than mocking it (workspace mocks-only-for-smoke rule)",
        "Per-class enforcement: Pathogen (R≥0.90, P≥0.95, F1≥0.92), Vaccine (R≥0.75, P≥0.85), "
        "Gene (R≥0.65, P≥0.95), Disease (search-only — no recall floor)",
        "Floors are LOWER BOUNDS, not target observed values — a build that misses any floor fails CI",
        "Source of truth: tests/integration/test_synonym_accuracy.py "
        "(behavior pinned by tests/unit/test_metrics_invariants.py)",
    ]
    _add_bullets(
        s,
        Inches(0.7),
        Inches(1.25),
        Inches(12.0),
        Inches(5.6),
        items,
        fontsize=15,
        line_spacing=1.3,
    )
    _add_footer(s, page, total)


def slide_accuracy_floors(prs, page, total):
    s = _new_slide(prs)
    _add_title_bar(
        s,
        "Accuracy Floors per Entity Class",
        "Slice baseline (60 rows) vs. full-corpus floor — both CI-enforced; bars are lower bounds, not observed values",
    )
    _add_image(
        s, FIG_DIR / "12_accuracy_thresholds.png", Inches(0.4), Inches(1.0), width=Inches(12.6)
    )
    _add_text_block(
        s,
        Inches(0.7),
        Inches(6.5),
        Inches(12.0),
        Inches(0.6),
        "A reading of 0.95 means the test enforces ≥95% — the actual run is at-or-above. "
        "Source: tests/integration/test_synonym_accuracy.py.",
        fontsize=12,
        color=SUBTLE,
    )
    _add_footer(s, page, total)


def slide_harmonization_stats(prs, page, total):
    s = _new_slide(prs)
    _add_title_bar(
        s,
        "Harmonization Statistics",
        "Source corpus rows · resolution-status confidence · 13,238 total rows verified 2026-05-05",
    )
    _add_image(
        s, FIG_DIR / "13_harmonization_stats.png", Inches(0.4), Inches(1.0), width=Inches(12.6)
    )
    _add_text_block(
        s,
        Inches(0.7),
        Inches(6.5),
        Inches(12.0),
        Inches(0.6),
        "Status taxonomy from synonym_dictionary/enums.py — every resolution result carries (status, confidence) "
        "to the caller, so downstream code can filter by quality.",
        fontsize=12,
        color=SUBTLE,
    )
    _add_footer(s, page, total)


def slide_session_bugs(prs, page, total):
    s = _new_slide(prs)
    _add_title_bar(
        s, "Session Outcome", "12 bugs / drifts uncovered, all fixed and pinned by tests"
    )
    _add_image(
        s, FIG_DIR / "09_session_bug_count.png", Inches(0.7), Inches(1.0), width=Inches(11.9)
    )
    _add_footer(s, page, total)


def slide_brutal_truth(prs, page, total):
    s = _new_slide(prs)
    _add_title_bar(
        s, "Things That Will Surprise You", "Silent-failure shapes you'll hit if you don't know"
    )
    items = [
        "Workflow.process(input) is fire-and-forget — use wait_for_cascade to await",
        "DirectLink defaults to auto_transfer=False — workflow loads, cascade no-ops",
        "Workflows need workflow-level input_data_units AND step-level data_units",
        "Trigger inputs are wrapped {unit_name: payload}; direct callers pass raw",
        "FAISS / sentence_transformers import order is load-bearing on macOS ARM",
        "Synthesis branch failures degrade gracefully; all-empty raises ValueError",
        "Globus Search is read-only at the ingest boundary — harvester is offline",
        "extra='forbid' is mandatory on every step config (silent typo defense)",
    ]
    _add_bullets(
        s,
        Inches(0.8),
        Inches(1.3),
        Inches(12.0),
        Inches(5.5),
        items,
        fontsize=18,
        line_spacing=1.25,
    )
    _add_footer(s, page, total)


def slide_failure_contract(prs, page, total):
    s = _new_slide(prs)
    _add_title_bar(s, "Failure Contract Per Branch", "Branch failures degrade; all-empty raises")
    rows = [
        ("Branch", "Failure", "Effect", True),
        (
            "Domain RAG",
            "Missing FAISS / corrupted bin",
            "rag_chunks=[], WARNING",
            False,
        ),
        (
            "VIOLIN / BV-BRC",
            "Missing CSV / wrong column",
            "both bundles=[], WARNING",
            False,
        ),
        ("PubMed", "Network / 5xx / timeout", "publications=[], WARNING", False),
        ("Globus Search", "SDK missing / network / bad UUID", "globus_results=[], WARNING", False),
        (
            "Synthesis LLM",
            "Endpoint down / model rejects",
            "ValueError → MCP {error}",
            False,
        ),
        (
            "ALL branches empty",
            "fail_on_empty_retrieval gate",
            "ValueError raised",
            False,
        ),
    ]
    table = s.shapes.add_table(
        rows=len(rows),
        cols=3,
        left=Inches(0.6),
        top=Inches(1.3),
        width=Inches(12.1),
        height=Inches(5.0),
    ).table
    col_widths = [Inches(2.6), Inches(4.7), Inches(4.8)]
    for ci, w in enumerate(col_widths):
        table.columns[ci].width = w
    for ri, (a, b, c, is_header) in enumerate(rows):
        for ci, txt in enumerate((a, b, c)):
            cell = table.cell(ri, ci)
            cell.text = txt
            cell.margin_left = Inches(0.15)
            cell.margin_right = Inches(0.15)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = Pt(15)
                    r.font.bold = is_header
                    r.font.color.rgb = TITLE_TXT if is_header else BODY_TXT
            if is_header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TITLE_BAR
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = (
                    RGBColor(0xF5, 0xF5, 0xF8) if ri % 2 == 1 else RGBColor(0xFF, 0xFF, 0xFF)
                )
    _add_footer(s, page, total)


def slide_close(prs, page, total):
    s = _new_slide(prs)
    _add_title_bar(s, "References", "Where to read more")
    items = [
        "docs/architecture.md — canonical end-to-end map (8 Mermaid diagrams + this content)",
        "CLAUDE.md — repo-local rules (Python venv, FAISS import order, MCP tool listing)",
        "../CLAUDE.md — workspace policy (cross-repo rules, mocks carve-out, harvester boundary)",
        "docs/figures/*.png — source images for this deck (rerun build_figures.py to regenerate)",
        "tests/integration/test_rag_e2e_workflow_yaml.py — the cascade test that found the silent-failure bugs",
    ]
    _add_bullets(
        s, Inches(0.8), Inches(1.4), Inches(12.0), Inches(5.0), items, fontsize=17, line_spacing=1.4
    )
    _add_text_block(
        s,
        Inches(0.8),
        Inches(6.5),
        Inches(11.7),
        Inches(0.5),
        "End — questions / discussion",
        fontsize=18,
        bold=True,
        color=ACCENT,
    )
    _add_footer(s, page, total)


def main() -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    builders = [
        slide_title,
        slide_agenda,
        slide_purpose,
        slide_two_lifecycles,
        slide_backend_pipeline,
        slide_user_facing_pipeline,
        slide_topology,
        slide_synthesis,
        slide_invocation,
        slide_trigger_cascade,
        slide_resolution,
        slide_ontologies,
        slide_mcp_tools,
        slide_tests,
        slide_accuracy_methodology,
        slide_accuracy_floors,
        slide_harmonization_stats,
        slide_session_bugs,
        slide_brutal_truth,
        slide_failure_contract,
        slide_close,
    ]
    total = len(builders)
    # Title slide doesn't take a count
    builders[0](prs)
    for page_num, b in enumerate(builders[1:], start=2):
        b(prs, page_num, total)
    prs.save(str(OUT_PATH))
    print(f"Slides written to {OUT_PATH}")
    print(f"  {OUT_PATH.stat().st_size // 1024} KB · {total} slides")


if __name__ == "__main__":
    main()
